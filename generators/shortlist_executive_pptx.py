"""Executive Shortlist PPTX generator (landscape, two slides per candidate).

Takes structured candidate data + client/role info, returns branded PPTX bytes.
Uses the landscape executive template with a data slide and a notes slide per candidate.
Font: Aptos throughout.
"""

import io
from copy import deepcopy
from datetime import datetime, date
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

from .shortlist_pptx import _strip_webextensions, _company_group_totals


TEMPLATE_PATH = Path(__file__).parent.parent / "templates" / "shortlist-executive-template.pptx"
PLACEHOLDER_PATH = Path(__file__).parent.parent / "assets" / "photo-placeholder.png"

FONT_NAME = "Aptos"

LOREM = (
    "Lorem ipsum dolor sit amet, consectetur adipiscing elit. "
    "Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua."
    "\x0b\x0b"
    "Ut enim ad minim veniam, quis nostrud exercitation ullamco laboris "
    "nisi ut aliquip ex ea commodo consequat."
    "\x0b\x0b"
    "Duis aute irure dolor in reprehenderit in voluptate velit esse cillum "
    "dolore eu fugiat nulla pariatur."
)


# ---------------------------------------------------------------------------
# Helpers (shared patterns from the standard generator)
# ---------------------------------------------------------------------------

def _calc_duration(start_str: str, end_str: str) -> str:
    """Calculate human-readable duration between two date strings."""
    if not start_str:
        return ""
    try:
        start = datetime.strptime(start_str, "%b %Y").date()
    except ValueError:
        return ""
    if not end_str or end_str.lower() == "present":
        end = date.today()
    else:
        try:
            end = datetime.strptime(end_str, "%b %Y").date()
        except ValueError:
            return ""
    total_months = (end.year - start.year) * 12 + end.month - start.month
    years, months = divmod(total_months, 12)
    if years == 0:
        return f"{months} month{'s' if months != 1 else ''}"
    elif months == 0:
        return f"{years} year{'s' if years != 1 else ''}"
    else:
        return f"{years} year{'s' if years != 1 else ''}, {months} month{'s' if months != 1 else ''}"


def _set_row_cell_text(row_elem, col_idx: int, text: str):
    """Set text in a cloned table row's cell via XML, preserving formatting."""
    cells = row_elem.findall(qn("a:tc"))
    if col_idx >= len(cells):
        return
    tc = cells[col_idx]
    txBody = tc.find(qn("a:txBody"))
    if txBody is None:
        return
    p = txBody.find(qn("a:p"))
    if p is None:
        return
    runs = p.findall(qn("a:r"))
    for extra in runs[1:]:
        p.remove(extra)
    r = p.find(qn("a:r"))
    if r is not None:
        t = r.find(qn("a:t"))
        if t is not None:
            t.text = text


def _set_detail_cell(cell, text: str):
    """Set a details table cell with proper line break handling."""
    p0 = cell.text_frame.paragraphs[0]
    p_elem = p0._p

    # Extract formatting
    rPr_template = None
    if p0.runs:
        rPr_elem = p0.runs[0]._r.find(qn("a:rPr"))
        if rPr_elem is not None:
            rPr_template = deepcopy(rPr_elem)
    if rPr_template is None:
        endRPr = p_elem.find(qn("a:endParaRPr"))
        if endRPr is not None:
            rPr_template = deepcopy(endRPr)
            rPr_template.tag = qn("a:rPr")
    if rPr_template is None:
        txBody_elem = p_elem.getparent()
        for other_p in txBody_elem.findall(qn("a:p")):
            for other_r in other_p.findall(qn("a:r")):
                rPr_found = other_r.find(qn("a:rPr"))
                if rPr_found is not None:
                    rPr_template = deepcopy(rPr_found)
                    break
            if rPr_template is not None:
                break

    # Clear existing runs and breaks
    for child in list(p_elem):
        if child.tag in (qn("a:r"), qn("a:br")):
            p_elem.remove(child)

    # Remove extra paragraphs
    txBody_elem = p_elem.getparent()
    for extra_p in list(txBody_elem.findall(qn("a:p")))[1:]:
        txBody_elem.remove(extra_p)

    # Insert new runs before endParaRPr
    endParaRPr = p_elem.find(qn("a:endParaRPr"))

    parts = text.replace("\n", "\x0b").split("\x0b")
    for j, part in enumerate(parts):
        if j > 0:
            br_elem = etree.Element(qn("a:br"))
            if endParaRPr is not None:
                endParaRPr.addprevious(br_elem)
            else:
                p_elem.append(br_elem)
        r_elem = etree.Element(qn("a:r"))
        if rPr_template is not None:
            r_elem.insert(0, deepcopy(rPr_template))
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = part
        if endParaRPr is not None:
            endParaRPr.addprevious(r_elem)
        else:
            p_elem.append(r_elem)


def _strip_italic_from_cell(cell):
    """Remove italic formatting from all runs in a cell."""
    for para in cell.text_frame.paragraphs:
        p_elem = para._p
        for r in p_elem.iter(qn("a:r")):
            rPr = r.find(qn("a:rPr"))
            if rPr is not None and rPr.get("i") is not None:
                del rPr.attrib["i"]
        for tag in ("a:defRPr", "a:endParaRPr"):
            rPr = p_elem.find(qn(tag))
            if rPr is not None and rPr.get("i") is not None:
                del rPr.attrib["i"]


def _replace_picture(slide, old_shape, new_image_bytes: bytes):
    """Replace a Picture shape's image, keeping position and size."""
    left, top = old_shape.left, old_shape.top
    width, height = old_shape.width, old_shape.height
    sp = old_shape._element
    sp.getparent().remove(sp)
    slide.shapes.add_picture(io.BytesIO(new_image_bytes), left, top, width, height)


def _clone_slide(prs: Presentation, slide_index: int) -> None:
    """Clone a slide and append it to the presentation."""
    import copy
    template_slide = prs.slides[slide_index]
    slide_layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    for shape in template_slide.shapes:
        el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.append(el)
    for ph in list(new_slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)


def _strip_theme_from_rPr(rPr, font_name: str):
    """Ensure rPr uses an explicit typeface with no theme override."""
    if rPr is None:
        return
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is not None:
            el.set("typeface", font_name)
            if "theme" in el.attrib:
                del el.attrib["theme"]
        elif tag == "a:latin":
            el = etree.SubElement(rPr, qn(tag))
            el.set("typeface", font_name)


def _fix_frame_fonts(tf, font_name: str):
    """Set font on all runs + paragraph defaults, stripping theme refs."""
    for para in tf.paragraphs:
        p_elem = para._p
        for run in para.runs:
            run.font.name = font_name
            rPr = run._r.find(qn("a:rPr"))
            _strip_theme_from_rPr(rPr, font_name)
        for tag in ("a:defRPr", "a:endParaRPr"):
            rPr = p_elem.find(qn(tag))
            if rPr is not None:
                _strip_theme_from_rPr(rPr, font_name)


def _strip_table_style(tbl_elem):
    """Remove table style reference so it can't re-impose theme fonts."""
    tblPr = tbl_elem.find(qn("a:tblPr"))
    if tblPr is not None:
        for ts in tblPr.findall(qn("a:tblStyle")):
            tblPr.remove(ts)


def _set_all_fonts(prs: Presentation, font_name: str):
    """Set all text to the specified font, stripping theme references.

    Theme font refs (e.g. theme='minor') on run properties and table styles
    prevent users from changing fonts in PowerPoint after export.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                _fix_frame_fonts(shape.text_frame, font_name)
            if shape.has_table:
                _strip_table_style(shape.table._tbl)
                for row in shape.table.rows:
                    for cell in row.cells:
                        _fix_frame_fonts(cell.text_frame, font_name)
            if hasattr(shape, "shapes"):
                for child in shape.shapes:
                    if child.has_text_frame:
                        _fix_frame_fonts(child.text_frame, font_name)


def _fill_data_slide(slide, cand: dict, placeholder_bytes: bytes | None):
    """Fill an executive shortlist data slide with candidate info."""
    career = [c for c in cand.get("career", []) if c.get("include", True)]
    photo_bytes = cand.get("photo") or placeholder_bytes

    for shape in slide.shapes:
        # Candidate name
        if shape.name == "TextBox 6":
            for para in shape.text_frame.paragraphs:
                for i, run in enumerate(para.runs):
                    run.text = cand["name"] if i == 0 else ""

        # Career history table
        elif shape.name == "Table 7":
            table = shape.table
            tbl = table._tbl

            template_row_xml = None
            if len(tbl.tr_lst) > 1:
                template_row_xml = deepcopy(tbl.tr_lst[1])
                for tc in template_row_xml.findall(qn("a:tc")):
                    for r in tc.iter(qn("a:r")):
                        rPr = r.find(qn("a:rPr"))
                        if rPr is not None:
                            rPr.set("sz", "700")

            while len(tbl.tr_lst) > 1:
                tbl.remove(tbl.tr_lst[-1])

            if template_row_xml is not None:
                group_totals = _company_group_totals(career)
                prev_company = None
                for idx, entry in enumerate(career):
                    new_row = deepcopy(template_row_xml)
                    company = entry.get("company", "")
                    title = entry.get("title", "")
                    start = entry.get("start_date", "")
                    end = entry.get("end_date", "")
                    duration = _calc_duration(start, end)

                    if company == prev_company:
                        display_company = ""
                    else:
                        display_company = company
                        total = group_totals.get(idx)
                        if total and company:
                            display_company = f"{company} ({total} total)"
                    prev_company = company

                    values = [display_company, title, start, end, duration]
                    for col_i, val in enumerate(values):
                        _set_row_cell_text(new_row, col_i, val)
                    tbl.append(new_row)

        # Details table
        elif shape.name == "Table 2":
            table = shape.table
            tbl_detail = table._tbl

            show_edu = cand.get("show_education", True)
            show_quals = cand.get("show_prof_quals", True)
            edu = cand.get("education", "")
            quals = cand.get("professional_qualifications", "")

            detail_data = [
                cand.get("notice_period", "") or "Not disclosed",
                cand.get("salary_expectation", "") or "Not disclosed",
                edu,
                quals,
            ]
            for row_i in range(min(4, len(table.rows))):
                _set_detail_cell(table.cell(row_i, 1), detail_data[row_i])

            # Strip italic from notice period row (row 0)
            if len(table.rows) >= 1:
                _strip_italic_from_cell(table.cell(0, 0))
                _strip_italic_from_cell(table.cell(0, 1))

            if not show_quals and len(tbl_detail.tr_lst) >= 4:
                tbl_detail.remove(tbl_detail.tr_lst[3])
            if not show_edu and len(tbl_detail.tr_lst) >= 3:
                tbl_detail.remove(tbl_detail.tr_lst[2])

        # Replace candidate photo
        elif shape.shape_type == 13 and photo_bytes:
            _replace_picture(slide, shape, photo_bytes)


def _fill_notes_slide(slide, cand: dict):
    """Fill an executive shortlist notes slide."""
    for shape in slide.shapes:
        if shape.name == "Rectangle: Rounded Corners 10":
            if shape.has_text_frame:
                tf = shape.text_frame
                notes_text = LOREM if cand.get("use_lorem", False) else cand.get("notes", "") or LOREM
                for p in tf.paragraphs:
                    for r in p.runs:
                        r.text = ""
                done = False
                for p in tf.paragraphs:
                    if p.runs and not done:
                        p.runs[0].text = notes_text
                        done = True


# ---------------------------------------------------------------------------
# Main generator
# ---------------------------------------------------------------------------

def generate_executive_shortlist(
    client_name: str,
    role_title: str,
    candidates: list[dict],
    prepared_by: str = "",
    prepared_date: str = "",
) -> bytes:
    """Generate a branded executive shortlist PPTX (landscape, two slides per candidate).

    Args:
        client_name: Client company name
        role_title: Role being recruited
        candidates: List of candidate dicts (same structure as standard generator)
        prepared_by: Consultant name for cover slide
        prepared_date: Date string for cover slide (e.g. "March 2026")

    Returns:
        PPTX file as bytes.
    """
    prs = Presentation(str(TEMPLATE_PATH))
    placeholder_bytes = PLACEHOLDER_PATH.read_bytes() if PLACEHOLDER_PATH.exists() else None

    if not prepared_date:
        prepared_date = datetime.now().strftime("%B %Y")

    # --- Cover slide (slide 0) ---
    slide0 = prs.slides[0]
    for shape in slide0.shapes:
        # Group 11 contains role title, "Shortlist", and client name
        if shape.name == "Group 11":
            for child in shape.shapes:
                if child.has_text_frame:
                    text = child.text_frame.text.strip()
                    if text == "Chief Executive Officer":
                        for p in child.text_frame.paragraphs:
                            for run in p.runs:
                                run.text = role_title
                    elif text == "Unico Group":
                        for p in child.text_frame.paragraphs:
                            for run in p.runs:
                                run.text = client_name

        # "object 3" has Prepared by and Date
        elif shape.name == "object 3" and shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                full_text = p.text
                if "Prepared by:" in full_text and p.runs:
                    # Last run contains the name
                    for run in p.runs:
                        if run.text.strip() and run.text.strip() != "Prepared by:":
                            run.text = f" {prepared_by}" if prepared_by else ""
                elif full_text.startswith("Date") and p.runs:
                    # Last run contains the date
                    for run in p.runs:
                        if run.text.strip() not in ("Date", ":"):
                            run.text = prepared_date

    # --- Ensure enough candidate slide pairs ---
    # Template has: slide 0 (cover), slide 1 (data), slide 2 (notes)
    # Each candidate needs a pair of slides
    needed = len(candidates)
    template_pairs = (len(prs.slides) - 1) // 2  # 1 pair in the template

    while template_pairs < needed:
        # Clone both the data slide and notes slide
        _clone_slide(prs, 1)  # clone data slide
        _clone_slide(prs, 2)  # clone notes slide
        template_pairs += 1

    # --- Fill each candidate's slide pair ---
    for cand_idx, cand in enumerate(candidates):
        data_slide_idx = 1 + cand_idx * 2
        notes_slide_idx = 2 + cand_idx * 2
        data_slide = prs.slides[data_slide_idx]
        notes_slide = prs.slides[notes_slide_idx]
        _fill_data_slide(data_slide, cand, placeholder_bytes)
        _fill_notes_slide(notes_slide, cand)

    # --- Remove extra candidate slide pairs from the end ---
    total_needed = 1 + needed * 2  # cover + pairs
    while len(prs.slides) > total_needed:
        rId = prs.slides._sldIdLst[-1].get(qn("r:id"))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])

    # --- Set all fonts to Aptos ---
    _set_all_fonts(prs, FONT_NAME)

    # --- Save ---
    buf = io.BytesIO()
    prs.save(buf)
    return _strip_webextensions(buf.getvalue())


def append_candidates(existing_pptx_bytes: bytes, candidates: list[dict]) -> bytes:
    """Append new candidates to an existing executive shortlist PPTX.

    Opens the existing file, clones the last candidate slide pair for each
    new candidate, fills them, and returns the updated PPTX bytes.
    Existing slides are not modified.
    """
    prs = Presentation(io.BytesIO(existing_pptx_bytes))
    placeholder_bytes = PLACEHOLDER_PATH.read_bytes() if PLACEHOLDER_PATH.exists() else None

    total = len(prs.slides)
    source_data_idx = total - 2
    source_notes_idx = total - 1

    for cand in candidates:
        _clone_slide(prs, source_data_idx)
        _clone_slide(prs, source_notes_idx)
        new_data_slide = prs.slides[len(prs.slides) - 2]
        new_notes_slide = prs.slides[len(prs.slides) - 1]
        _fill_data_slide(new_data_slide, cand, placeholder_bytes)
        _fill_notes_slide(new_notes_slide, cand)

    _set_all_fonts(prs, FONT_NAME)

    buf = io.BytesIO()
    prs.save(buf)
    return _strip_webextensions(buf.getvalue())
