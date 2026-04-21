"""Shortlist PPTX generator.

Takes structured candidate data + client/role info, returns branded PPTX bytes.
Uses the Infinitas shortlist template (Tate's edited version).
Font: Aptos throughout.
"""

import io
import re
import zipfile
from copy import deepcopy
from datetime import datetime, date
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu


def _strip_webextensions(pptx_bytes: bytes) -> bytes:
    """Remove Office web add-in (taskpane) parts from a PPTX.

    The Infinitas shortlist template carries a WA200010001 Office add-in that
    auto-launches a taskpane and blocks text editing in tables. Strip it from
    every generated file.
    """
    webext_rel_re = re.compile(
        r'<Relationship Id="[^"]*" Type="http://schemas\.microsoft\.com/office/2011/relationships/webextensiontaskpanes" Target="[^"]*"/>')
    webext_ct_re = re.compile(r'<Override PartName="/ppt/webextensions/[^"]*"[^/]*/>')

    src = zipfile.ZipFile(io.BytesIO(pptx_bytes), 'r')
    out_buf = io.BytesIO()
    with zipfile.ZipFile(out_buf, 'w', zipfile.ZIP_DEFLATED) as dst:
        for item in src.infolist():
            name = item.filename
            if name.startswith('ppt/webextensions/'):
                continue
            data = src.read(name)
            if name == '_rels/.rels':
                data = webext_rel_re.sub('', data.decode('utf-8')).encode('utf-8')
            elif name == '[Content_Types].xml':
                data = webext_ct_re.sub('', data.decode('utf-8')).encode('utf-8')
            dst.writestr(item, data)
    src.close()
    return out_buf.getvalue()


TEMPLATE_PATH = Path(__file__).parent.parent / "templates" / "shortlist-template.pptx"
PLACEHOLDER_PATH = Path(__file__).parent.parent / "assets" / "photo-placeholder.png"

FONT_NAME = "Aptos"

LOREM = (
    "Lorem ipsum dolor sit amet, consectetur adipiscing elit. "
    "Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua."
    "\x0b\x0b"
    "Ut enim ad minim veniam, quis nostrud exercitation ullamco laboris "
    "nisi ut aliquip ex ea commodo consequat."
    "\x0b\x0b"
    "Duis aute irure dolor in reprehenderit in voluptate velit esse cillum "
    "dolore eu fugiat nulla pariatur."
)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _calc_duration(start_str: str, end_str: str) -> str:
    """Calculate human-readable duration between two date strings."""
    if not start_str:
        return ""
    try:
        start = datetime.strptime(start_str, "%b %Y").date()
    except ValueError:
        return ""
    if not end_str or end_str.lower() == "present":
        end = date.today()
    else:
        try:
            end = datetime.strptime(end_str, "%b %Y").date()
        except ValueError:
            return ""
    total_months = (end.year - start.year) * 12 + end.month - start.month
    years, months = divmod(total_months, 12)
    if years == 0:
        return f"{months} month{'s' if months != 1 else ''}"
    elif months == 0:
        return f"{years} year{'s' if years != 1 else ''}"
    else:
        return f"{years} year{'s' if years != 1 else ''}, {months} month{'s' if months != 1 else ''}"


def _company_group_totals(career: list[dict]) -> dict[int, str]:
    """For each run of consecutive career entries at the same company, return
    a map from the group's first-row index to a "X years, Y months" span
    (earliest start -> latest end across all roles in the group).

    Single-role groups are excluded — their row already shows the duration.
    """
    totals: dict[int, str] = {}
    i = 0
    n = len(career)
    while i < n:
        company = career[i].get("company", "")
        group_end = i
        while group_end + 1 < n and career[group_end + 1].get("company", "") == company:
            group_end += 1
        if group_end > i:
            group = career[i : group_end + 1]
            # Earliest start across the group
            earliest = None
            for e in group:
                s = e.get("start_date", "")
                if not s:
                    continue
                try:
                    d = datetime.strptime(s, "%b %Y").date()
                except ValueError:
                    continue
                if earliest is None or d < earliest[0]:
                    earliest = (d, s)
            # Latest end — "Present" dominates
            has_current = any(
                not e.get("end_date") or e.get("end_date", "").lower() == "present"
                for e in group
            )
            if has_current:
                latest_str = "Present"
            else:
                latest = None
                for e in group:
                    ed = e.get("end_date", "")
                    if not ed:
                        continue
                    try:
                        d = datetime.strptime(ed, "%b %Y").date()
                    except ValueError:
                        continue
                    if latest is None or d > latest[0]:
                        latest = (d, ed)
                latest_str = latest[1] if latest else ""
            if earliest and latest_str:
                totals[i] = _calc_duration(earliest[1], latest_str)
        i = group_end + 1
    return totals


def _set_row_cell_text(row_elem, col_idx: int, text: str):
    """Set text in a cloned table row's cell via XML, preserving formatting."""
    cells = row_elem.findall(qn("a:tc"))
    if col_idx >= len(cells):
        return
    tc = cells[col_idx]
    txBody = tc.find(qn("a:txBody"))
    if txBody is None:
        return
    p = txBody.find(qn("a:p"))
    if p is None:
        return
    runs = p.findall(qn("a:r"))
    for extra in runs[1:]:
        p.remove(extra)
    r = p.find(qn("a:r"))
    if r is not None:
        t = r.find(qn("a:t"))
        if t is not None:
            t.text = text


def _set_row_cell_multi_paragraph(row_elem, col_idx: int, lines: list[str]):
    """Set a cell to multiple paragraphs, cloning formatting from the first one.

    Used for the Company column when there are multiple roles at a company —
    line 1 is the company name, line 2 is the total tenure.
    """
    cells = row_elem.findall(qn("a:tc"))
    if col_idx >= len(cells):
        return
    tc = cells[col_idx]
    txBody = tc.find(qn("a:txBody"))
    if txBody is None:
        return
    existing = txBody.findall(qn("a:p"))
    if not existing:
        return
    template_p = deepcopy(existing[0])
    for p in existing:
        txBody.remove(p)
    for line in lines:
        new_p = deepcopy(template_p)
        runs = new_p.findall(qn("a:r"))
        for extra in runs[1:]:
            new_p.remove(extra)
        r = new_p.find(qn("a:r"))
        if r is not None:
            t = r.find(qn("a:t"))
            if t is not None:
                t.text = line
        txBody.append(new_p)


def _strip_cell_fill(tc):
    """Remove any solid/no/grad/pat fills from a table cell's tcPr."""
    tcPr = tc.find(qn("a:tcPr"))
    if tcPr is None:
        return
    for tag in ("a:solidFill", "a:noFill", "a:gradFill", "a:blipFill", "a:pattFill"):
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)


def _apply_cell_fill_scheme(tc, scheme_val: str):
    """Apply a solid schemeClr fill (e.g. 'bg1' or 'bg2') to a table cell."""
    tcPr = tc.find(qn("a:tcPr"))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn("a:tcPr"))
    _strip_cell_fill(tc)
    solid = etree.SubElement(tcPr, qn("a:solidFill"))
    sch = etree.SubElement(solid, qn("a:schemeClr"))
    sch.set("val", scheme_val)


def _set_column_widths_emu(tbl_elem, widths_emu: list[int]):
    """Override each <a:gridCol w=...> in the table's <a:tblGrid>."""
    grid = tbl_elem.find(qn("a:tblGrid"))
    if grid is None:
        return
    cols = grid.findall(qn("a:gridCol"))
    for col, w in zip(cols, widths_emu):
        col.set("w", str(w))


def _set_detail_cell(cell, text: str):
    """Set a details table cell with proper line break handling."""
    p0 = cell.text_frame.paragraphs[0]
    p_elem = p0._p

    # Extract formatting
    rPr_template = None
    if p0.runs:
        rPr_elem = p0.runs[0]._r.find(qn("a:rPr"))
        if rPr_elem is not None:
            rPr_template = deepcopy(rPr_elem)
    if rPr_template is None:
        endRPr = p_elem.find(qn("a:endParaRPr"))
        if endRPr is not None:
            rPr_template = deepcopy(endRPr)
            rPr_template.tag = qn("a:rPr")
    if rPr_template is None:
        txBody_elem = p_elem.getparent()
        for other_p in txBody_elem.findall(qn("a:p")):
            for other_r in other_p.findall(qn("a:r")):
                rPr_found = other_r.find(qn("a:rPr"))
                if rPr_found is not None:
                    rPr_template = deepcopy(rPr_found)
                    break
            if rPr_template is not None:
                break

    # Clear existing runs and breaks
    for child in list(p_elem):
        if child.tag in (qn("a:r"), qn("a:br")):
            p_elem.remove(child)

    # Remove extra paragraphs
    txBody_elem = p_elem.getparent()
    for extra_p in list(txBody_elem.findall(qn("a:p")))[1:]:
        txBody_elem.remove(extra_p)

    # Insert new runs before endParaRPr
    endParaRPr = p_elem.find(qn("a:endParaRPr"))

    parts = text.replace("\n", "\x0b").split("\x0b")
    for j, part in enumerate(parts):
        if j > 0:
            br_elem = etree.Element(qn("a:br"))
            if endParaRPr is not None:
                endParaRPr.addprevious(br_elem)
            else:
                p_elem.append(br_elem)
        r_elem = etree.Element(qn("a:r"))
        if rPr_template is not None:
            r_elem.insert(0, deepcopy(rPr_template))
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = part
        if endParaRPr is not None:
            endParaRPr.addprevious(r_elem)
        else:
            p_elem.append(r_elem)


def _strip_italic_from_cell(cell):
    """Remove italic formatting from all runs in a cell."""
    for para in cell.text_frame.paragraphs:
        p_elem = para._p
        for r in p_elem.iter(qn("a:r")):
            rPr = r.find(qn("a:rPr"))
            if rPr is not None and rPr.get("i") is not None:
                del rPr.attrib["i"]
        for tag in ("a:defRPr", "a:endParaRPr"):
            rPr = p_elem.find(qn(tag))
            if rPr is not None and rPr.get("i") is not None:
                del rPr.attrib["i"]


def _replace_picture(slide, old_shape, new_image_bytes: bytes):
    """Replace a Picture shape's image, keeping position and size."""
    left, top = old_shape.left, old_shape.top
    width, height = old_shape.width, old_shape.height
    sp = old_shape._element
    sp.getparent().remove(sp)
    slide.shapes.add_picture(io.BytesIO(new_image_bytes), left, top, width, height)


def _clone_slide(prs: Presentation, slide_index: int) -> None:
    """Clone a slide and append it to the presentation."""
    import copy
    template_slide = prs.slides[slide_index]
    slide_layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    for shape in template_slide.shapes:
        el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.append(el)
    for ph in list(new_slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)


def _strip_theme_from_rPr(rPr, font_name: str):
    """Ensure rPr uses an explicit typeface with no theme override."""
    if rPr is None:
        return
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is not None:
            el.set("typeface", font_name)
            if "theme" in el.attrib:
                del el.attrib["theme"]
        elif tag == "a:latin":
            el = etree.SubElement(rPr, qn(tag))
            el.set("typeface", font_name)


def _fix_frame_fonts(tf, font_name: str):
    """Set font on all runs + paragraph defaults, stripping theme refs."""
    for para in tf.paragraphs:
        p_elem = para._p
        for run in para.runs:
            run.font.name = font_name
            rPr = run._r.find(qn("a:rPr"))
            _strip_theme_from_rPr(rPr, font_name)
        for tag in ("a:defRPr", "a:endParaRPr"):
            rPr = p_elem.find(qn(tag))
            if rPr is not None:
                _strip_theme_from_rPr(rPr, font_name)


def _strip_table_style(tbl_elem):
    """Remove table style reference so it can't re-impose theme fonts."""
    tblPr = tbl_elem.find(qn("a:tblPr"))
    if tblPr is not None:
        for ts in tblPr.findall(qn("a:tblStyle")):
            tblPr.remove(ts)


def _set_all_fonts(prs: Presentation, font_name: str):
    """Set all text to the specified font, stripping theme references.

    Theme font refs (e.g. theme='minor') on run properties and table styles
    prevent users from changing fonts in PowerPoint after export.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                _fix_frame_fonts(shape.text_frame, font_name)
            if shape.has_table:
                _strip_table_style(shape.table._tbl)
                for row in shape.table.rows:
                    for cell in row.cells:
                        _fix_frame_fonts(cell.text_frame, font_name)
            if hasattr(shape, "shapes"):
                for child in shape.shapes:
                    if child.has_text_frame:
                        _fix_frame_fonts(child.text_frame, font_name)


def _fill_candidate_slide(slide, cand: dict, placeholder_bytes: bytes | None):
    """Fill a single candidate slide with data."""
    career = [c for c in cand.get("career", []) if c.get("include", True)]
    photo_bytes = cand.get("photo") or placeholder_bytes

    for shape in slide.shapes:
        # Candidate name — centre over photo
        if shape.name == "TextBox 6":
            for para in shape.text_frame.paragraphs:
                for i, run in enumerate(para.runs):
                    run.text = cand["name"] if i == 0 else ""
            photo_shape = None
            for s in slide.shapes:
                if s.shape_type == 13:
                    photo_shape = s
                    break
            if photo_shape is not None:
                photo_cx = photo_shape.left + photo_shape.width // 2
                shape.left = photo_cx - shape.width // 2

        # Career history table
        elif shape.name == "Table 7":
            table = shape.table
            tbl = table._tbl

            template_row_xml = None
            if len(tbl.tr_lst) > 1:
                template_row_xml = deepcopy(tbl.tr_lst[1])
                for tc in template_row_xml.findall(qn("a:tc")):
                    for p in tc.iter(qn("a:p")):
                        pPr = p.find(qn("a:pPr"))
                        if pPr is None:
                            pPr = etree.SubElement(p, qn("a:pPr"))
                            p.insert(0, pPr)
                        defRPr = pPr.find(qn("a:defRPr"))
                        if defRPr is None:
                            defRPr = etree.SubElement(pPr, qn("a:defRPr"))
                        defRPr.set("sz", "900")
                    for r in tc.iter(qn("a:r")):
                        rPr = r.find(qn("a:rPr"))
                        if rPr is not None and "sz" in rPr.attrib:
                            del rPr.attrib["sz"]

            while len(tbl.tr_lst) > 1:
                tbl.remove(tbl.tr_lst[-1])

            if template_row_xml is not None:
                # Column widths that Tate settled on manually: narrow dates,
                # wider Company column to hold "Company / X years total" on two lines.
                _EMU = 914400
                _CAREER_WIDTHS = (1.78, 2.34, 0.91, 0.83, 1.28)
                _set_column_widths_emu(tbl, [int(w * _EMU) for w in _CAREER_WIDTHS])

                # Build groups of consecutive roles at the same company.
                groups: list[tuple[str, list[int]]] = []
                i = 0
                while i < len(career):
                    comp = career[i].get("company", "")
                    j = i
                    while j + 1 < len(career) and career[j + 1].get("company", "") == comp:
                        j += 1
                    groups.append((comp, list(range(i, j + 1))))
                    i = j + 1

                group_totals = _company_group_totals(career)

                for group_idx, (company, indices) in enumerate(groups):
                    # Alternate banding per group (not per row), so a merged
                    # multi-row group reads as a single band.
                    scheme = "bg1" if group_idx % 2 == 0 else "bg2"
                    total = group_totals.get(indices[0])

                    for pos, row_idx in enumerate(indices):
                        entry = career[row_idx]
                        new_row = deepcopy(template_row_xml)
                        title = entry.get("title", "")
                        start = entry.get("start_date", "")
                        end = entry.get("end_date", "")
                        duration = _calc_duration(start, end)

                        cells = new_row.findall(qn("a:tc"))

                        # Company column: merged across multi-role groups,
                        # with "Company" / "X years total" stacked on two lines.
                        if pos == 0:
                            if len(indices) > 1 and total:
                                _set_row_cell_multi_paragraph(
                                    new_row, 0, [company, f"{total} total"]
                                )
                                cells[0].set("rowSpan", str(len(indices)))
                            else:
                                _set_row_cell_text(new_row, 0, company)
                        else:
                            _set_row_cell_text(new_row, 0, "")
                            cells[0].set("vMerge", "1")

                        _set_row_cell_text(new_row, 1, title)
                        _set_row_cell_text(new_row, 2, start)
                        _set_row_cell_text(new_row, 3, end)
                        _set_row_cell_text(new_row, 4, duration)

                        # Apply band fill to every cell in this row (including
                        # the merged Company cell on its top row; continuation
                        # cells are invisible anyway but kept consistent).
                        for c in cells:
                            _apply_cell_fill_scheme(c, scheme)

                        tbl.append(new_row)

        # Details table
        elif shape.name == "Table 2":
            table = shape.table
            tbl_detail = table._tbl

            show_edu = cand.get("show_education", True)
            show_quals = cand.get("show_prof_quals", True)
            edu = cand.get("education", "")
            quals = cand.get("professional_qualifications", "")

            detail_data = [
                cand.get("notice_period", "") or "Not disclosed",
                cand.get("salary_expectation", "") or "Not disclosed",
                edu,
                quals,
            ]
            for row_i in range(min(4, len(table.rows))):
                _set_detail_cell(table.cell(row_i, 1), detail_data[row_i])

            # Strip italic from notice period row (row 0)
            if len(table.rows) >= 1:
                _strip_italic_from_cell(table.cell(0, 0))
                _strip_italic_from_cell(table.cell(0, 1))

            if not show_quals and len(tbl_detail.tr_lst) >= 4:
                tbl_detail.remove(tbl_detail.tr_lst[3])
            if not show_edu and len(tbl_detail.tr_lst) >= 3:
                tbl_detail.remove(tbl_detail.tr_lst[2])

        # Notes
        elif shape.name in ("Rectangle: Rounded Corners 8", "Rectangle: Rounded Corners 10"):
            if shape.has_text_frame:
                tf = shape.text_frame
                notes_text = LOREM if cand.get("use_lorem", False) else cand.get("notes", "") or LOREM
                for p in tf.paragraphs:
                    for r in p.runs:
                        r.text = ""
                done = False
                for p in tf.paragraphs:
                    if p.runs and not done:
                        p.runs[0].text = notes_text
                        done = True

        # Replace candidate photo
        elif shape.shape_type == 13 and photo_bytes:
            _replace_picture(slide, shape, photo_bytes)


# ---------------------------------------------------------------------------
# Main generator
# ---------------------------------------------------------------------------

def generate_shortlist(
    client_name: str,
    role_title: str,
    candidates: list[dict],
) -> bytes:
    """Generate a branded shortlist PPTX.

    Args:
        client_name: Client company name
        role_title: Role being recruited
        candidates: List of candidate dicts with:
            - name: str
            - career: list[dict] with company, title, start_date, end_date, include
            - education: str
            - professional_qualifications: str
            - show_education: bool (True to include Education row)
            - show_prof_quals: bool (True to include Prof Quals row)
            - notice_period: str
            - salary_expectation: str
            - notes: str
            - use_lorem: bool
            - photo: bytes or None

    Returns:
        PPTX file as bytes.
    """
    prs = Presentation(str(TEMPLATE_PATH))
    placeholder_bytes = PLACEHOLDER_PATH.read_bytes() if PLACEHOLDER_PATH.exists() else None

    # --- Cover slide (slide 0) ---
    slide0 = prs.slides[0]
    for shape in slide0.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if "CHIEF EXECUTIVE OFFICER" in text:
                for run in para.runs:
                    run.text = role_title
            elif text == "Unico Group":
                for run in para.runs:
                    run.text = client_name

    # --- Ensure enough candidate slides ---
    template_candidate_count = len(prs.slides) - 1
    needed = len(candidates)
    while template_candidate_count < needed:
        _clone_slide(prs, 1)
        template_candidate_count += 1

    # --- Fill each candidate slide ---
    for cand_idx, cand in enumerate(candidates):
        slide = prs.slides[cand_idx + 1]
        _fill_candidate_slide(slide, cand, placeholder_bytes)

    # --- Remove extra candidate slides ---
    while len(prs.slides) > needed + 1:
        rId = prs.slides._sldIdLst[-1].get(qn("r:id"))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])

    # --- Set all fonts to Aptos ---
    _set_all_fonts(prs, FONT_NAME)

    # --- Save ---
    buf = io.BytesIO()
    prs.save(buf)
    return _strip_webextensions(buf.getvalue())


def append_candidates(existing_pptx_bytes: bytes, candidates: list[dict]) -> bytes:
    """Append new candidates to an existing standard shortlist PPTX.

    Opens the existing file, clones the last candidate slide for each new
    candidate, fills it, and returns the updated PPTX bytes.
    Existing slides are not modified.
    """
    prs = Presentation(io.BytesIO(existing_pptx_bytes))
    placeholder_bytes = PLACEHOLDER_PATH.read_bytes() if PLACEHOLDER_PATH.exists() else None

    source_idx = len(prs.slides) - 1

    for cand in candidates:
        _clone_slide(prs, source_idx)
        new_slide = prs.slides[len(prs.slides) - 1]
        _fill_candidate_slide(new_slide, cand, placeholder_bytes)

    _set_all_fonts(prs, FONT_NAME)

    buf = io.BytesIO()
    prs.save(buf)
    return _strip_webextensions(buf.getvalue())
